from django.shortcuts import render

# Create your views here.
from django.shortcuts import render, redirect, get_object_or_404
from .forms import PresentationForm
from .models import PresentationData
from pptx import Presentation
from django.http import HttpResponse, FileResponse
import pandas as pd
import io

from django.shortcuts import render, get_object_or_404
from django.http import JsonResponse, FileResponse
from .forms import PresentationForm
from .models import PresentationData
from django.contrib.auth.decorators import user_passes_test
from pptx import Presentation
import os
from django.http import FileResponse
from django.shortcuts import render
from .forms import PresentationForm
from .models import PresentationData

from django.http import FileResponse
from django.shortcuts import render
from .forms import PresentationForm
from .models import PresentationData
from datetime import date
import os
from pptx import Presentation



from django.contrib.auth.decorators import user_passes_test

from django.contrib.auth.decorators import user_passes_test

def superuser_required(view_func):
    def check_user(u):
        return u.is_active and (u.is_superuser or u.is_authenticated)
    return user_passes_test(check_user)(view_func)



from django.contrib.auth import logout
from django.shortcuts import redirect

def user_logout(request):
    logout(request)
    return redirect('login')


CONSTITUENCIES_BY_DISTRICT={
  "Srikakulam": ["Ichchapuram", "Palasa", "Tekkali", "Pathapatnam", "Srikakulam", "Amadalavalasa", "Etcherla", "Narasannapeta"],
  "Vizianagaram": ["Rajam", "Palakonda", "Kurupam", "Parvathipuram", "Salur", "Bobbili", "Cheepurupalli", "Gajapathinagaram", "Nellimarla", "Vizianagaram"],
  "Parvathipuram Manyam": ["Kurupam", "Parvathipuram", "Salur", "Araku Valley"],
  "Visakhapatnam": ["Srungavarapukota", "Bhimili", "Visakhapatnam East", "Visakhapatnam South", "Visakhapatnam North", "Visakhapatnam West", "Gajuwaka", "Chodavaram", "Elamanchili", "Pendurthi"],
  "Anakapalli": ["Madugula", "Anakapalle", "Pendurthi", "Payakaraopet", "Tuni", "Narsipatnam", "Chodavaram", "Pendurthi"],
  "Alluri Sitharama Raju": ["Araku Valley", "Paderu", "Rampachodravaram"],
  "Kakinada": ["Tuni", "Prathipadu", "Pithapuram", "Kakinada Rural", "Kakinada City", "Ramachandrapuram", "Mummidivaram"],
  "East Godavari": ["Anaparthy", "Kakinada City", "Ramachandrapuram", "Mummidivaram", "Razole", "Kothapeta", "Mandapeta"],
  "Konaseema": ["Amalapuram", "Razole", "Mummidivaram", "Kothapeta", "Ramachandrapuram", "Mummidivaram", "Kothapeta"],
  "West Godavari": ["Achanta", "Palakollu", "Bhimavaram", "Undi", "Tanuku", "Tadepalligudem", "Unguturu"],
  "Eluru": ["Unguturu", "Denduluru", "Polavaram", "Chintalapudi", "Tiruvuru", "Nuzvid", "Kaikalur"],
  "NTR": ["Nuzvid", "Gannavaram", "Kaikalur", "Pedana", "Machilipatnam", "Avanigadda", "Pamarru"],
  "Krishna": ["Gannavaram", "Gudivada", "Kaikalur", "Pedana", "Machilipatnam", "Avanigadda", "Pamarru"],
  "Guntur": ["Pedakurapadu", "Tadikonda", "Mangalagiri", "Ponnuru", "Vemuru", "Repalle", "Tenali"],
  "Palnadu": ["Chilakaluripet", "Narasaraopet", "Sattenapalle", "Vinukonda", "Gurazala", "Macherla", "Yerragondapalem"],
  "Bapatla": ["Darsi", "Parchur", "Addanki", "Chirala", "Santhanuthalapadu", "Ongole", "Markapuram"],
  "Prakasam": ["Yerragondapalem", "Darsi", "Parchur", "Addanki", "Chirala", "Santhanuthalapadu", "Ongole", "Kandukur", "Markapuram"],
  "Nellore": ["Kondapi", "Kavali", "Atmakur", "Nellore City", "Nellore Rural", "Udayagiri"],
  "Tirupati": ["Kavali", "Gudur", "Sullurpeta", "Nagari", "Tirupati", "Srikalahasti"],
  "Chittoor": ["Punganur", "Chandragiri", "Tirupati", "Satyavedu", "Nagari", "Chittoor", "Puthalapattu", "Palamaner", "Kuppam"],
  "Annamayya": ["Rajampet", "Kodur", "Rayachoti", "Thamballapalle", "Pileru", "Madanapalle"],
  "YSR Kadapa": ["Badvel", "Kadapa", "Rayachoti", "Pulivendla", "Kamalapuram", "Jammalamadugu", "Proddatur", "Mydukur"],
  "Nandyal": ["Allagadda", "Nandyal", "Nandikotkur", "Banaganapalle", "Dhone", "Pattikonda"],
  "Kurnool": ["Kurnool", "Pattikonda", "Adoni", "Alur", "Kodumur", "Yemmiganur", "Mantralayam"],
  "Ananthapuramu": ["Rayadurg", "Uravakonda", "Guntakal", "Tadipatri", "Singanamala", "Anantapur Urban", "Kalyandurg", "Hindupur"],
  "Sri Sathya Sai": ["Madakasira", "Hindupur", "Penukonda", "Puttaparthi", "Dharmavaram", "Kadiri"]
}
import os
from django.shortcuts import render
from django.http import FileResponse, HttpResponseBadRequest
from .models import PresentationData
from pptx import Presentation
from datetime import datetime
from io import BytesIO

@superuser_required
def create_presentation(request):
    if request.method == "POST":
        # Extract fields from POST
        date_str = request.POST.get('date', '')
        toname = request.POST.get('toname', '')
        phno = request.POST.get('phno', '')
        aadharno = request.POST.get('aadharno', '')
        ap = request.POST.get('ap', '')
        aptdas = request.POST.get('aptdas', '')
        ap_district = request.POST.get('ap_district', '')
        ap_constitution = request.POST.get('ap_constitution', '')
        address = request.POST.get('address', '')
        pincode=request.POST.get('pincode','')
        committee=request.POST.get('committee','')
        if not all([date_str, toname, phno, aadharno, ap, aptdas, ap_district, ap_constitution, address,pincode,committee]):
            return HttpResponseBadRequest("All fields are required.")

        try:
            date_obj = datetime.strptime(date_str, '%Y-%m-%d').date()
        except ValueError:
            return HttpResponseBadRequest("Invalid date format.")
        temp=aptdas
        ap_parts = ap.split(' ', 2)
        ap1 = ap_parts[0]+" "+ap_parts[1]
        ap2 = ap_parts[2] if len(ap_parts) > 2 else ''
        aptdas=aptdas+" "+"-"+' '+committee
        data = PresentationData(
            date=date_obj,
            toname=toname,
            phno=phno,
            aadharno=aadharno,
            ap1=ap1,
            ap2=ap2,
            aptdas=temp,
            address=address,
            ap_district=ap_district,
            ap_constitution=ap_constitution,
            pincode=pincode,
            committee=committee
        )
        data.save()

        # Load your PPTX template
        template_path = "Presentation1.pptx"
        prs = Presentation(template_path)

        # Map placeholders to data, adjust keys to match your PPTX placeholders
        replace_map = {
            "date": data.date.strftime('%d-%m-%Y'),
            "toname": data.toname,
            "phno": data.phno,
            "aadharno": data.aadharno,
            "ap1": data.ap1,
            "ap2": data.ap2,
            "aptdas": aptdas,
            "address": data.address,
            "ap_district": data.ap_district,
            "ap_constitution": data.ap_constitution,
        }

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            for key, val in replace_map.items():
                                if key in run.text:
                                    run.text = run.text.replace(key, val)

        # Save pptx to BytesIO stream (no file on disk)
        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        # Prepare response to send pptx file for download
        response = FileResponse(
            pptx_io,
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        )
        # Suggested filename; adjust as needed
        filename = f"{data.toname}_Appointment_Letter.pptx"
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response

    # GET request renders form template
    context = {
        'today': date.today().strftime('%Y-%m-%d'),
    }
    return render(request, "presentation_form.html",context)



from django.shortcuts import render
from .models import PresentationData

@superuser_required
def track_history(request):
    entries = PresentationData.objects.all().order_by('-created_at')
    return render(request, 'history.html', {
        'entries': entries,
        'search_term': '',
        'not_found': False,
    })

@superuser_required
def export_excel(request):
    entries = PresentationData.objects.all().values(
        'entry_id', 'date', 'toname', 'phno', 'aadharno', 'ap1', 'ap2', 'aptdas', 'address', 'ap_district', 'ap_constitution','pincode', 'created_at'
    )
    import pandas as pd
    import io
    from django.http import HttpResponse

    df = pd.DataFrame(entries)
    
    # Add serial number column starting from 1
    df.insert(0, 'S.No', range(1, len(df) + 1))
    df.rename(columns={'aptdas': 'Designation'}, inplace=True)
    # Convert datetimes to strings (or remove tz info)
    if 'created_at' in df.columns:
        df['created_at'] = df['created_at'].apply(
            lambda dt: dt.strftime('%d-%m-%Y %H:%M') if hasattr(dt, 'strftime') else str(dt)
        )
    if 'date' in df.columns:
        df['date'] = df['date'].apply(
            lambda d: d.strftime('%d-%m-%Y') if hasattr(d, 'strftime') else str(d)
        )

    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine='xlsxwriter') as writer:
        df.to_excel(writer, index=False, sheet_name='Presentations')
    buffer.seek(0)

    response = HttpResponse(buffer, content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    response['Content-Disposition'] = 'attachment; filename="Members_Data.xlsx"'
    return response


from django.http import FileResponse
from django.shortcuts import get_object_or_404
from .models import PresentationData
import os
from pptx import Presentation



from django.http import FileResponse
from django.shortcuts import get_object_or_404
from .models import PresentationData
import os
from pptx import Presentation

from django.shortcuts import get_object_or_404
from django.http import FileResponse
from pptx import Presentation
from io import BytesIO
@superuser_required


def history_download_ppt(request, pk):
    entry = get_object_or_404(PresentationData, pk=pk)

    # Prepare file (no special filename sanitization needed for in-memory download)
    template_path = "Presentation1.pptx"

    # Build aptdas logic exactly as in your create view (to include committee)
    aptdas_with_committee = entry.aptdas + " - " + entry.committee if entry.committee else entry.aptdas
    temp=entry.ap1.split()
    entry.ap1=temp[0]+" "+temp[1]
    entry.ap2=temp[2]+" "+temp[3]
    replace_map = {
        "date": entry.date.strftime('%d-%m-%Y'),
        "toname": entry.toname,
        "phno": entry.phno,
        "aadharno": entry.aadharno,
        "ap1": entry.ap1,
        "ap2": entry.ap2,
        "aptdas": aptdas_with_committee,
        "address": entry.address,
        "ap_district": entry.ap_district,
        "ap_constitution": entry.ap_constitution,
    }
    entry.save()

    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for key, val in replace_map.items():
                            if key in run.text:
                                run.text = run.text.replace(key, str(val))

    # Save pptx to BytesIO (not disk)
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)

    # Suggested filename; uses name and Entry id
    filename = f"{entry.toname}_Appointment_Letter.pptx"
    response = FileResponse(
        pptx_io,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    return response




def get_constituencies(request):
    district = request.GET.get('district', '')
    constituencies = CONSTITUENCIES_BY_DISTRICT.get(district, [])
    return JsonResponse({'constituencies': constituencies})
def history_edit(request, pk):
    entry = get_object_or_404(PresentationData, pk=pk)

    if request.method == 'POST':
        # Collect updated form data from POST request
        date = request.POST.get('date')
        toname = request.POST.get('toname')
        phno = request.POST.get('phno')
        aadharno = request.POST.get('aadharno')
        ap = request.POST.get('ap')
        aptdas = request.POST.get('aptdas')
        ap_district = request.POST.get('ap_district')
        ap_constitution = request.POST.get('ap_constitution')
        address = request.POST.get('address')
        pincode = request.POST.get('pincode')
        committee=request.POST.get('committee')
        # Update entry object fields
        entry.date = date
        entry.toname = toname
        entry.phno = phno
        entry.aadharno = aadharno
        entry.ap1 = ap   # Assuming field ap1 matches this
        entry.aptdas = aptdas
        entry.ap_district = ap_district
        entry.ap_constitution = ap_constitution
        entry.address = address
        entry.pincode = pincode
        entry.committee=committee
        entry.save()
        return redirect('track_history')  # Adjust redirect name as needed

    else:
        # GET request, pre-populate form fields with existing entry data
        context = {
            'entry': entry,
            'today': entry.date,  # Assuming you want to default date input to existing date
        }
        return render(request, 'history_edit.html', context)
    
from django.shortcuts import get_object_or_404, redirect
from django.views.decorators.http import require_POST
from .models import PresentationData

@require_POST
def history_delete(request, pk):
    entry = get_object_or_404(PresentationData, pk=pk)
    entry.delete()
    return redirect('track_history')  # Make sure you have a URL with name='history'
